"""Generate unit-wise PPTX from the sample template with expanded content slides."""

from __future__ import annotations

import re
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

from .models import PlannedUnit, Subject, UnitContentPackage


def _set_shape_text(shape, text: str, font_size: int | None = None) -> None:
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    if font_size:
        p.font.size = Pt(font_size)


def _set_placeholder_body(slide, text: str) -> None:
    for shape in slide.shapes:
        if shape.name.startswith("Text Placeholder") and shape.has_text_frame:
            _set_shape_text(shape, text)
            return
    for shape in slide.shapes:
        if shape.has_text_frame and shape.shape_type == 14:
            _set_shape_text(shape, text)
            return


def _bullet_lines(items: list[str], max_items: int = 8) -> str:
    lines = []
    for item in items[:max_items]:
        item = re.sub(r"\s+", " ", item).strip()
        if item:
            lines.append(f"• {item}")
    return "\n".join(lines) if lines else "• (Content to be added)"


def _add_content_textbox(slide, text: str) -> None:
    for shape in list(slide.shapes):
        if shape.shape_type == 17 and "TextBox" in shape.name:
            sp = shape.element
            sp.getparent().remove(sp)
    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(7.5), Inches(4.8))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text[:4000]
    p.font.size = Pt(14)


def _duplicate_slide(prs: Presentation, index: int):
    template = prs.slides[index]
    layout = template.slide_layout
    new_slide = prs.slides.add_slide(layout)
    for sh in list(new_slide.shapes):
        el = sh.element
        el.getparent().remove(el)
    for shape in template.shapes:
        newel = deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(newel, "p:extLst")
    return new_slide


def _move_slide(prs: Presentation, old_index: int, new_index: int) -> None:
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    el = slides[old_index]
    xml_slides.remove(el)
    xml_slides.insert(new_index, el)


def _fill_content_slide(slide, title: str, bullets: list[str], notes: str = "") -> None:
    for shape in slide.shapes:
        if shape.name.startswith("Google Shape;98") and shape.has_text_frame:
            _set_shape_text(shape, title[:100], 24)
    body = _bullet_lines(bullets, 6)
    if notes:
        body += f"\n\n{notes[:500]}"
    _add_content_textbox(slide, body)


def generate_unit_ppt(
    template_path: Path,
    output_path: Path,
    *,
    program_name: str,
    course_code: str,
    subject: Subject,
    unit: PlannedUnit,
    pack: UnitContentPackage | None = None,
    delivered_by: str = "",
) -> Path:
    from .content_engine import generate_unit_content

    if pack is None:
        pack = generate_unit_content(subject, unit, max(1, len(subject.syllabus_units) or subject.credits))

    prs = Presentation(str(template_path))
    topics_list = [s.title for s in pack.topic_slides] or unit.topics

    # Slide 1 — Title
    for shape in prs.slides[0].shapes:
        if shape.name.startswith("Google Shape") and shape.has_text_frame:
            tf = shape.text_frame
            tf.clear()
            p0 = tf.paragraphs[0]
            p0.text = unit.title
            p0.font.size = Pt(28)
            p1 = tf.add_paragraph()
            p1.text = f"Delivered By: {delivered_by}".strip()
            p1.font.size = Pt(18)
        if "Footer" in shape.name and shape.has_text_frame:
            _set_shape_text(
                shape,
                f"Program Name: {program_name}    |    Course Code: {course_code}",
            )

    # Slide 2 — Syllabus / intro
    syllabus_text = pack.introduction or subject.description or f"Overview of {subject.name}"
    _set_placeholder_body(prs.slides[1], syllabus_text[:3500])

    # Slide 3 — Topics
    _set_placeholder_body(prs.slides[2], _bullet_lines(topics_list, 12))

    # Slide 4 — Objectives
    for shape in prs.slides[3].shapes:
        if shape.name.startswith("Google Shape;90") and shape.has_text_frame:
            _set_shape_text(
                shape,
                "By the end of this session, you will be able to:\n\n"
                + _bullet_lines(pack.learning_objectives, 6),
                16,
            )

    # Expand content slides (insert before summary at index 5)
    summary_idx = 5
    content_slides = pack.topic_slides
    _fill_content_slide(
        prs.slides[4],
        content_slides[0].title if content_slides else "Introduction",
        content_slides[0].bullets if content_slides else ["Unit content"],
        content_slides[0].speaker_notes if content_slides else "",
    )

    for i in range(1, len(content_slides)):
        dup_idx = len(prs.slides)
        _duplicate_slide(prs, 4)
        _move_slide(prs, dup_idx, summary_idx + i)
        _fill_content_slide(
            prs.slides[summary_idx + i],
            content_slides[i].title,
            content_slides[i].bullets,
            content_slides[i].speaker_notes,
        )

    # Case study slide before summary (optional extra)
    if pack.case_study and len(prs.slides) < pack.target_slide_count:
        dup_idx = len(prs.slides)
        _duplicate_slide(prs, 4)
        insert_at = len(prs.slides) - 3
        _move_slide(prs, dup_idx, insert_at)
        _fill_content_slide(prs.slides[insert_at], "Case Study / Example", [pack.case_study[:300]])

    # Summary — always last content before closing
    final_summary_idx = len(prs.slides) - 2
    summary_points = pack.keywords[:4] or unit.topics[-4:] or [pack.summary[:200]]
    _set_placeholder_body(
        prs.slides[final_summary_idx],
        "Key takeaways:\n\n" + pack.summary[:1500] + "\n\n" + _bullet_lines(summary_points, 5),
    )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


def generate_all_unit_ppts(
    template_path: Path,
    output_dir: Path,
    *,
    program_name: str,
    course_code: str,
    subject: Subject,
    units: list[PlannedUnit],
    packs: dict[str, UnitContentPackage] | None = None,
    delivered_by: str = "",
    use_ai: bool | None = None,
) -> list[Path]:
    from .content_engine import generate_unit_content

    paths: list[Path] = []
    packs = packs or {}
    unit_count = len(units)
    safe = re.sub(r"[^\w\-]+", "_", subject.name)[:60].strip("_")

    for unit in units:
        key = unit.label
        if key not in packs:
            packs[key] = generate_unit_content(
                subject, unit, unit_count, use_ai=use_ai
            )
        pack = packs[key]
        fname = f"{safe}_Unit_{unit.label}.pptx"
        out = output_dir / fname
        generate_unit_ppt(
            template_path,
            out,
            program_name=program_name,
            course_code=course_code,
            subject=subject,
            unit=unit,
            pack=pack,
            delivered_by=delivered_by,
        )
        paths.append(out)
    return paths
